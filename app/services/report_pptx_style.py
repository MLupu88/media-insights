"""MSL brand styling constants and shape-building helpers for PPTX report
generation (Phase 6C).

Derived directly from the confirmed web-app brand tokens
(`tailwind.config.js`) — there is no separate template *file*; every slide
applies these constants explicitly through python-pptx's shape/text-frame
API, the same "one shared source of styling constants, reused everywhere"
discipline the Jinja macros already apply to HTML. See the Phase 6C plan's
"Template finding" section for why no file from `review H1/` is used.
"""

from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

from app.services.report_contract import MAX_PPTX_LABEL_CHARS

# --- MSL brand tokens (tailwind.config.js) -----------------------------------
INK = RGBColor(0x11, 0x11, 0x11)
PAPER = RGBColor(0xFA, 0xF9, 0xF6)
LINE = RGBColor(0xE4, 0xE1, 0xDA)
ACCENT = RGBColor(0xB3, 0x49, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x6B, 0x6B, 0x6B)  # readable secondary text on `PAPER`
ZEBRA_FILL = RGBColor(0xF2, 0xF0, 0xEC)  # light tint between PAPER and LINE, alternating body rows

DISPLAY_FONT = "Newsreader"
SANS_FONT = "Inter"

# 16:9
SLIDE_WIDTH = Emu(12192000)
SLIDE_HEIGHT = Emu(6858000)

LOGO_PATH = "app/static/brand/msl-the-practice-logo-cropped.png"
LOGO_WIDTH = Inches(1.4)
LOGO_HEIGHT = Emu(int(LOGO_WIDTH * 60 / 420))  # preserve the 420x60 logo aspect ratio

MARGIN = Inches(0.5)
CONTENT_TOP = Inches(1.15)

TITLE_FONT_SIZE = Pt(30)
SECTION_FONT_SIZE = Pt(22)
SUBTITLE_FONT_SIZE = Pt(12)
BODY_FONT_SIZE = Pt(11)
CAPTION_FONT_SIZE = Pt(9)

# A table's height is derived from its actual row count, never stretched to
# an arbitrary caller-supplied height — the previous approach let
# PowerPoint distribute a large fixed height evenly across a sparse table's
# rows, turning the solid-fill header row into an oversized black block on
# any slide with few rows (the exact bug this fixes).
TABLE_ROW_HEIGHT = Inches(0.32)
TABLE_MAX_HEIGHT = Inches(3.6)


def truncate_label(text: str | None, max_chars: int = MAX_PPTX_LABEL_CHARS) -> str:
    """Deterministic slide-only label truncation — the Excel export never
    truncates a label; this exists purely for chart/table legibility.
    """
    if not text:
        return ""
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 1].rstrip() + "…"


def add_background(slide) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER


def add_logo(slide) -> None:
    from pptx.util import Length

    left: Length = SLIDE_WIDTH - LOGO_WIDTH - MARGIN
    slide.shapes.add_picture(LOGO_PATH, left, MARGIN, width=LOGO_WIDTH, height=LOGO_HEIGHT)


def add_blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # "Blank"
    add_background(slide)
    add_logo(slide)
    return slide


def add_title_text(slide, text: str, subtitle: str | None = None) -> None:
    """Large centered title, used only on the title slide."""
    box = slide.shapes.add_textbox(MARGIN, Inches(2.9), SLIDE_WIDTH - 2 * MARGIN, Inches(1.6))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = 2  # PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = DISPLAY_FONT
    run.font.size = TITLE_FONT_SIZE
    run.font.bold = True
    run.font.color.rgb = INK

    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = 2
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.name = SANS_FONT
        run2.font.size = SUBTITLE_FONT_SIZE
        run2.font.color.rgb = MUTED


def add_accent_bar(slide, top=Inches(0.28)) -> None:
    """A thin MSL-accent-colored rule under the section header — the one
    deliberate accent-color brand touch on every content slide, matching
    the terracotta accent used for emphasis/interactive states throughout
    the web app.
    """
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, top, Inches(0.6), Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_section_title(slide, text: str, caption: str | None = None) -> None:
    """Section header used on every content slide, top-left, with an
    optional small caption directly under it (e.g. "Top 10 of 34 shown").
    """
    add_accent_bar(slide)
    box = slide.shapes.add_textbox(MARGIN, Inches(0.35), SLIDE_WIDTH - 2 * MARGIN, Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = DISPLAY_FONT
    run.font.size = SECTION_FONT_SIZE
    run.font.bold = True
    run.font.color.rgb = INK

    if caption:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = caption
        run2.font.name = SANS_FONT
        run2.font.size = CAPTION_FONT_SIZE
        run2.font.italic = True
        run2.font.color.rgb = MUTED


def add_label_badge(slide, text: str) -> None:
    """A small accent-filled badge, top-right — used to visually
    distinguish "Interpretation" from "Recommendation" on Key Findings
    slides at a glance, not just by reading the text.
    """
    width = Inches(0.25) + Pt(len(text) * 6)
    box = slide.shapes.add_textbox(SLIDE_WIDTH - MARGIN - width, Inches(0.4), width, Inches(0.4))
    box.fill.solid()
    box.fill.fore_color.rgb = ACCENT
    tf = box.text_frame
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(3)
    tf.margin_bottom = Pt(3)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text.upper()
    run.font.name = SANS_FONT
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = WHITE


def add_body_text(slide, text: str, top, height=Inches(1.0)) -> None:
    box = slide.shapes.add_textbox(MARGIN, top, SLIDE_WIDTH - 2 * MARGIN, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = SANS_FONT
    run.font.size = BODY_FONT_SIZE
    run.font.color.rgb = INK


def add_no_data_message(slide, top, message: str = "No data available for the selected filters.") -> None:
    """Compact fallback for an empty section — a table is never rendered
    with zero data rows, which would otherwise show as a bare, disproportionately
    tall header block once PowerPoint distributes a caller height across it.
    """
    add_body_text(slide, message, top, Inches(0.5))


def add_table(slide, headers: list[str], rows: list[list[str]], top) -> "Length":
    """Renders a header row (black fill, white text) followed by body rows
    (white/zebra fill, dark text). Row height is fixed per row and derived
    from the actual row count — never stretched to an arbitrary height —
    and callers must guard against `rows == []` themselves by rendering
    `add_no_data_message` instead of calling this with zero rows.

    Returns the table's actual bottom Y position (EMU) so callers can place
    any following content (e.g. a chart) directly under it instead of
    assuming a fixed slide layout.
    """
    from pptx.util import Emu as _Emu

    n_rows = len(rows) + 1
    n_cols = len(headers)
    width = SLIDE_WIDTH - 2 * MARGIN

    row_height = TABLE_ROW_HEIGHT
    total_height = row_height * n_rows
    if total_height > TABLE_MAX_HEIGHT:
        row_height = _Emu(int(TABLE_MAX_HEIGHT / n_rows))
        total_height = row_height * n_rows

    graphic_frame = slide.shapes.add_table(n_rows, n_cols, MARGIN, top, width, total_height)
    table = graphic_frame.table
    for row in table.rows:
        row.height = row_height

    for col_index, header in enumerate(headers):
        cell = table.cell(0, col_index)
        cell.text = header
        cell.margin_top = Pt(2)
        cell.margin_bottom = Pt(2)
        cell.vertical_anchor = 3  # MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = INK
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = SANS_FONT
                run.font.size = Pt(10)
                run.font.bold = True
                run.font.color.rgb = WHITE

    for row_index, row_values in enumerate(rows, start=1):
        row_fill = WHITE if row_index % 2 == 1 else ZEBRA_FILL
        for col_index, value in enumerate(row_values):
            cell = table.cell(row_index, col_index)
            cell.text = str(value)
            cell.margin_top = Pt(2)
            cell.margin_bottom = Pt(2)
            cell.vertical_anchor = 3  # MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_fill
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = SANS_FONT
                    run.font.size = BODY_FONT_SIZE
                    run.font.color.rgb = INK

    return top + total_height

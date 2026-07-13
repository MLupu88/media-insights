"""Builds the PowerPoint export (Phase 6C) from a `ProjectReportData`/
`ComparisonReportData` instance (`app/services/report_data.py`).

One function per slide type; two orchestrators (`build_project_pptx`,
`build_comparison_pptx`). Every slide applies the MSL style constants from
`report_pptx_style.py`. Charts are native python-pptx chart objects, never
rasterized images. Every ranked section is capped at `REPORT_TOP_N` with an
explicit "(Top N of M shown)" caption whenever the underlying data has
more; the full ranking always lives in the matching Excel export instead.
"""

from io import BytesIO

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from app.services import report_pptx_style as style
from app.services.report_contract import (
    MAX_INSIGHT_TEXT_CHARS_PER_SLIDE,
    MAX_PPTX_BYTES,
    NO_DATA_MESSAGE,
    REPORT_TOP_N,
)
from app.services.report_data import ComparisonReportData, ProjectReportData, ReportTooLargeError

CHART_HEIGHT = Inches(2.3)
CHART_GAP = Inches(0.25)


# --- formatting helpers --------------------------------------------------------


def _pct(value) -> str:
    return f"{value:.1f}%" if value is not None else "—"


def _pp(value) -> str:
    if value is None:
        return "—"
    sign = "+" if value > 0 else ""
    return f"{sign}{value:.1f}pp"


def _num(value) -> str:
    if value is None:
        return "—"
    return f"{value:,.0f}"


def _clip_insight_text(text: str) -> str:
    if len(text) <= MAX_INSIGHT_TEXT_CHARS_PER_SLIDE:
        return text
    return text[:MAX_INSIGHT_TEXT_CHARS_PER_SLIDE].rstrip() + "… (see Excel export for full text)"


def _cap_caption(noun: str, full_list: list, top_n: int = REPORT_TOP_N) -> str | None:
    if len(full_list) > top_n:
        return f"Top {top_n} of {len(full_list)} {noun} shown — full ranking in the Excel export"
    return None


def _bar_chart(slide, categories: list[str], series_name: str, values: list[float], top) -> None:
    if not categories:
        return
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(series_name, values)
    slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.5),
        top,
        style.SLIDE_WIDTH - Inches(1.0),
        CHART_HEIGHT,
        chart_data,
    )


def _add_ranked_table(slide, headers: list[str], rows: list[list[str]], top):
    """Renders the section's table, or a compact "No data available"
    message in its place when `rows` is empty — never an empty
    multi-row/black table. Returns the bottom Y position to place any
    following content (e.g. a chart), or `None` when no table was drawn.
    """
    if not rows:
        style.add_no_data_message(slide, top)
        return None
    return style.add_table(slide, headers, rows, top)


def _insight_slides(prs, insights, heading: str) -> None:
    for insight in insights:
        slide = style.add_blank_slide(prs)
        style.add_section_title(slide, heading, insight.narrative_type.replace("_", " "))
        style.add_label_badge(slide, insight.label)
        style.add_body_text(slide, insight.title, style.CONTENT_TOP, Inches(0.6))
        style.add_body_text(slide, _clip_insight_text(insight.narrative), Inches(1.9), Inches(3.0))
        citation_bits = [
            bit
            for bit in (insight.related_brand, insight.related_publication)
            if bit
        ]
        if citation_bits:
            style.add_body_text(slide, "Related: " + ", ".join(citation_bits), Inches(5.2), Inches(0.5))


def _methodology_slide(prs, metadata, extra_lines: list[str] | None = None) -> None:
    slide = style.add_blank_slide(prs)
    style.add_section_title(slide, "Methodology & Filters")
    coverage = metadata.article_detail_coverage
    insight_coverage = metadata.insight_coverage
    lines = [
        f"Scope: {metadata.scope_label}",
        f"Filters: {metadata.filters_label}",
        f"Generated: {metadata.generated_at.strftime('%Y-%m-%d %H:%M UTC')}",
        metadata.population_definition,
        metadata.ai_methodology_note,
        metadata.chat_exclusion_note,
        f"Article detail: showing {coverage.shown_count:,} of {coverage.total_count:,} eligible articles"
        + (" (truncated)" if coverage.truncated else ""),
        f"Insights: {insight_coverage.available_count} valid insight(s) available, "
        f"{insight_coverage.included_count} included, "
        f"{insight_coverage.excluded_causal_count} excluded (unsupported causal language)",
    ]
    if extra_lines:
        lines.extend(extra_lines)
    style.add_body_text(slide, "\n".join(lines), style.CONTENT_TOP, Inches(5.5))


def _check_size(prs) -> bytes:
    buffer = BytesIO()
    prs.save(buffer)
    data = buffer.getvalue()
    if len(data) > MAX_PPTX_BYTES:
        raise ReportTooLargeError(
            f"Generated PowerPoint exceeded the maximum supported size "
            f"({MAX_PPTX_BYTES:,} bytes)."
        )
    return data


# --- project report slides -----------------------------------------------------


def _project_title_slide(prs, data: ProjectReportData) -> None:
    slide = style.add_blank_slide(prs)
    style.add_title_text(
        slide,
        data.project_name,
        f"Media Coverage Report — {data.project_quarter}\n"
        f"Generated {data.metadata.generated_at.strftime('%Y-%m-%d %H:%M UTC')}",
    )


def _no_data_slide(prs) -> None:
    slide = style.add_blank_slide(prs)
    style.add_section_title(slide, "No Data Available")
    style.add_body_text(slide, NO_DATA_MESSAGE, style.CONTENT_TOP)


def _project_executive_summary_slide(prs, data: ProjectReportData) -> None:
    slide = style.add_blank_slide(prs)
    style.add_section_title(slide, "Executive Summary")
    kpis = data.analytics["kpis"]
    headers = ["Metric", "Value"]
    rows = [
        ["Unique valid articles", _num(kpis["unique_valid_articles"])],
        ["Total reach", _num(kpis["total_reach"])],
        ["Average reach", _num(kpis["average_reach"])],
        ["Publications", _num(kpis["publication_count"])],
        ["Classified articles", _num(kpis["unique_classified_articles"])],
    ]
    bottom = style.add_table(slide, headers, rows, style.CONTENT_TOP)

    insight = next((i for i in data.insights if i.narrative_type == "executive_summary"), None)
    if insight:
        text_top = bottom + CHART_GAP
        style.add_body_text(
            slide, f"{insight.label}: {_clip_insight_text(insight.narrative)}", text_top, Inches(2.6)
        )


def _project_brand_slide(prs, data: ProjectReportData) -> None:
    slide = style.add_blank_slide(prs)
    full = data.analytics["brands"]["by_volume"]
    style.add_section_title(slide, "Brand & Competitor Performance", _cap_caption("brands", full))
    top_rows = full[:REPORT_TOP_N]
    headers = ["Brand", "Articles", "SOV %", "Reach Share %"]
    rows = [
        [style.truncate_label(r["brand"]), _num(r["article_count"]), _pct(r["sov_pct"]), _pct(r["reach_share_pct"])]
        for r in top_rows
    ]
    bottom = _add_ranked_table(slide, headers, rows, style.CONTENT_TOP)
    if bottom is not None:
        _bar_chart(
            slide, [style.truncate_label(r["brand"]) for r in top_rows], "SOV %",
            [r["sov_pct"] for r in top_rows], bottom + CHART_GAP,
        )


def _project_topic_slide(prs, data: ProjectReportData) -> None:
    slide = style.add_blank_slide(prs)
    full = data.analytics["topics"]["primary_topic_distribution"]
    style.add_section_title(slide, "Topic & Category Mix", _cap_caption("topics", full))
    top_rows = full[:REPORT_TOP_N]
    headers = ["Topic", "Count", "%"]
    rows = [[style.truncate_label(r["value"]), _num(r["count"]), _pct(r["pct"])] for r in top_rows]
    bottom = _add_ranked_table(slide, headers, rows, style.CONTENT_TOP)
    if bottom is not None:
        _bar_chart(
            slide, [style.truncate_label(r["value"]) for r in top_rows], "Count",
            [r["count"] for r in top_rows], bottom + CHART_GAP,
        )


def _project_sentiment_slide(prs, data: ProjectReportData) -> None:
    slide = style.add_blank_slide(prs)
    sentiment = data.analytics["sentiment"]
    full = sentiment["overall_distribution"]
    style.add_section_title(slide, "Sentiment & Brand Role", _cap_caption("sentiment values", full))
    top_rows = full[:REPORT_TOP_N]
    headers = ["Sentiment", "Count", "%"]
    rows = [[style.truncate_label(r["value"]), _num(r["count"]), _pct(r["pct"])] for r in top_rows]
    bottom = _add_ranked_table(slide, headers, rows, style.CONTENT_TOP)
    if bottom is not None:
        _bar_chart(slide, [r["value"] for r in top_rows], "Count", [r["count"] for r in top_rows], bottom + CHART_GAP)


def _project_publications_slide(prs, data: ProjectReportData) -> None:
    slide = style.add_blank_slide(prs)
    full = data.analytics["publications_and_stories"]["publications_by_volume"]
    style.add_section_title(slide, "Publications", _cap_caption("publications", full))
    top_rows = full[:REPORT_TOP_N]
    headers = ["Publication", "Articles", "Volume %", "Reach %"]
    rows = [
        [style.truncate_label(r["publication"]), _num(r["article_count"]), _pct(r["volume_pct"]), _pct(r["reach_pct"])]
        for r in top_rows
    ]
    bottom = _add_ranked_table(slide, headers, rows, style.CONTENT_TOP)
    if bottom is not None:
        _bar_chart(
            slide,
            [style.truncate_label(r["publication"]) for r in top_rows],
            "Articles",
            [r["article_count"] for r in top_rows],
            bottom + CHART_GAP,
        )


def _project_stories_slide(prs, data: ProjectReportData) -> None:
    slide = style.add_blank_slide(prs)
    full = data.analytics["publications_and_stories"]["stories_by_volume"]
    style.add_section_title(slide, "Story Clusters", _cap_caption("stories", full))
    top_rows = full[:REPORT_TOP_N]
    headers = ["Story", "Articles", "Total Reach"]
    rows = [
        [style.truncate_label(r["story_key"]), _num(r["article_count"]), _num(r["total_reach"])]
        for r in top_rows
    ]
    bottom = _add_ranked_table(slide, headers, rows, style.CONTENT_TOP)
    if bottom is not None:
        _bar_chart(
            slide,
            [style.truncate_label(r["story_key"]) for r in top_rows],
            "Articles",
            [r["article_count"] for r in top_rows],
            bottom + CHART_GAP,
        )


def build_project_pptx(data: ProjectReportData) -> bytes:
    prs = Presentation()
    prs.slide_width = style.SLIDE_WIDTH
    prs.slide_height = style.SLIDE_HEIGHT

    _project_title_slide(prs, data)

    if data.analytics["kpis"]["unique_valid_articles"] == 0:
        _no_data_slide(prs)
        _methodology_slide(prs, data.metadata)
        return _check_size(prs)

    _project_executive_summary_slide(prs, data)
    _project_brand_slide(prs, data)
    _project_topic_slide(prs, data)
    _project_sentiment_slide(prs, data)
    _project_publications_slide(prs, data)
    _project_stories_slide(prs, data)
    _insight_slides(prs, data.insights, "Key Findings")
    _methodology_slide(prs, data.metadata)

    return _check_size(prs)


# --- comparison report slides ---------------------------------------------------


def _comparison_title_slide(prs, data: ComparisonReportData) -> None:
    slide = style.add_blank_slide(prs)
    style.add_title_text(
        slide,
        "Comparison Report",
        f"{data.baseline_label} vs {data.comparison_label}\n"
        f"Generated {data.metadata.generated_at.strftime('%Y-%m-%d %H:%M UTC')}",
    )


def _comparison_executive_summary_slide(prs, data: ComparisonReportData) -> None:
    slide = style.add_blank_slide(prs)
    style.add_section_title(slide, "Executive Summary")
    kpi_deltas = data.comparison["deltas"]["kpis"]
    headers = ["Metric", "Baseline", "Comparison", "Abs. Delta", "% Delta"]
    rows = []
    for key, label in (
        ("unique_valid_articles", "Unique valid articles"),
        ("total_reach", "Total reach"),
        ("average_reach", "Average reach"),
        ("publication_count", "Publications"),
    ):
        d = kpi_deltas[key]
        rows.append(
            [label, _num(d["baseline"]), _num(d["comparison"]), _num(d["absolute_delta"]),
             _pct(d["percentage_delta"]) if d["percentage_delta"] is not None else "—"]
        )
    bottom = style.add_table(slide, headers, rows, style.CONTENT_TOP)

    insight = next(
        (i for i in data.insights if i.narrative_type == "comparison_executive_summary"), None
    )
    if insight:
        text_top = bottom + CHART_GAP
        style.add_body_text(
            slide, f"{insight.label}: {_clip_insight_text(insight.narrative)}", text_top, Inches(2.6)
        )


def _comparison_brand_slide(prs, data: ComparisonReportData) -> None:
    slide = style.add_blank_slide(prs)
    full = data.comparison["deltas"]["brands"]
    style.add_section_title(slide, "Brand & SOV Movement", _cap_caption("brands", full))
    top_rows = full[:REPORT_TOP_N]
    headers = ["Brand", "Baseline SOV %", "Comparison SOV %", "SOV Delta"]
    rows = []
    for r in top_rows:
        tag = " (new)" if r["is_new_entrant"] else " (dropped)" if r["is_dropout"] else ""
        rows.append(
            [style.truncate_label(r["brand"]) + tag, _pct(r["baseline_sov_pct"]),
             _pct(r["comparison_sov_pct"]), _pp(r["sov_delta_pp"])]
        )
    bottom = _add_ranked_table(slide, headers, rows, style.CONTENT_TOP)
    if bottom is not None:
        _bar_chart(
            slide,
            [style.truncate_label(r["brand"]) for r in top_rows],
            "SOV Delta (pp)",
            [r["sov_delta_pp"] for r in top_rows],
            bottom + CHART_GAP,
        )


def _comparison_topic_slide(prs, data: ComparisonReportData) -> None:
    slide = style.add_blank_slide(prs)
    full = data.comparison["deltas"]["topics"]["rows"]
    style.add_section_title(slide, "Topic & Category Shifts", _cap_caption("topics", full))
    top_rows = full[:REPORT_TOP_N]
    headers = ["Topic", "Baseline %", "Comparison %", "Delta (pp)"]
    rows = [
        [style.truncate_label(r["value"]), _pct(r["baseline_pct"]), _pct(r["comparison_pct"]), _pp(r["pct_delta_pp"])]
        for r in top_rows
    ]
    _add_ranked_table(slide, headers, rows, style.CONTENT_TOP)


def _comparison_sentiment_slide(prs, data: ComparisonReportData) -> None:
    slide = style.add_blank_slide(prs)
    full = data.comparison["deltas"]["sentiment"]["rows"]
    style.add_section_title(slide, "Sentiment & Brand-Role Change", _cap_caption("sentiment values", full))
    top_rows = full[:REPORT_TOP_N]
    headers = ["Sentiment", "Baseline %", "Comparison %", "Delta (pp)"]
    rows = [
        [style.truncate_label(r["value"]), _pct(r["baseline_pct"]), _pct(r["comparison_pct"]), _pp(r["pct_delta_pp"])]
        for r in top_rows
    ]
    _add_ranked_table(slide, headers, rows, style.CONTENT_TOP)


def _comparison_publications_slide(prs, data: ComparisonReportData) -> None:
    slide = style.add_blank_slide(prs)
    full = data.comparison["deltas"]["publications_by_volume"]
    style.add_section_title(slide, "Publication Movement", _cap_caption("publications", full))
    top_rows = full[:REPORT_TOP_N]
    headers = ["Publication", "Baseline Rank", "Comparison Rank", "Movement"]
    rows = []
    for r in top_rows:
        tag = "New" if r["is_new_entrant"] else "Dropped" if r["is_dropout"] else str(r["rank_change"] or "—")
        rows.append(
            [style.truncate_label(r["publication"]), r["baseline_rank"] or "—", r["comparison_rank"] or "—", tag]
        )
    _add_ranked_table(slide, headers, rows, style.CONTENT_TOP)


def _comparison_stories_slide(prs, data: ComparisonReportData) -> None:
    slide = style.add_blank_slide(prs)
    full = data.comparison["deltas"]["stories_by_volume"]
    style.add_section_title(slide, "Story Movement", _cap_caption("stories", full))
    top_rows = full[:REPORT_TOP_N]
    headers = ["Story", "Baseline Rank", "Comparison Rank", "Movement"]
    rows = []
    for r in top_rows:
        tag = "New" if r["is_new_entrant"] else "Dropped" if r["is_dropout"] else str(r["rank_change"] or "—")
        rows.append(
            [style.truncate_label(r["story_key"]), r["baseline_rank"] or "—", r["comparison_rank"] or "—", tag]
        )
    _add_ranked_table(slide, headers, rows, style.CONTENT_TOP)


def _comparison_volatility_slide(prs, data: ComparisonReportData) -> None:
    slide = style.add_blank_slide(prs)
    style.add_section_title(slide, "Volatility & Concentration")
    volatility = data.comparison["volatility"]
    concentration = data.comparison["deltas"]["publication_concentration"]
    headers = ["Metric", "Value"]
    rows = [
        ["Publications: avg rank change", _num(volatility["publications"]["avg_rank_change"])],
        ["Publications: new entrants", _num(volatility["publications"]["entrants_count"])],
        ["Publications: dropouts", _num(volatility["publications"]["dropouts_count"])],
        ["Stories: avg rank change", _num(volatility["stories"]["avg_rank_change"])],
        ["Stories: new entrants", _num(volatility["stories"]["entrants_count"])],
        ["Stories: dropouts", _num(volatility["stories"]["dropouts_count"])],
        ["Top 3 publication volume concentration delta", _pp(concentration["top3_volume_pct"]["delta_pp"])],
    ]
    style.add_table(slide, headers, rows, style.CONTENT_TOP)


def build_comparison_pptx(data: ComparisonReportData) -> bytes:
    prs = Presentation()
    prs.slide_width = style.SLIDE_WIDTH
    prs.slide_height = style.SLIDE_HEIGHT

    _comparison_title_slide(prs, data)

    baseline_empty = data.comparison["baseline"]["kpis"]["unique_valid_articles"] == 0
    comparison_empty = data.comparison["comparison"]["kpis"]["unique_valid_articles"] == 0
    if baseline_empty and comparison_empty:
        _no_data_slide(prs)
        _methodology_slide(prs, data.metadata)
        return _check_size(prs)

    _comparison_executive_summary_slide(prs, data)
    _comparison_brand_slide(prs, data)
    _comparison_topic_slide(prs, data)
    _comparison_sentiment_slide(prs, data)
    _comparison_publications_slide(prs, data)
    _comparison_stories_slide(prs, data)
    _comparison_volatility_slide(prs, data)
    _insight_slides(prs, data.insights, "Key Findings")
    _methodology_slide(prs, data.metadata)

    return _check_size(prs)

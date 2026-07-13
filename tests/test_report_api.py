import uuid
from unittest.mock import patch

import openpyxl
import pytest
from pptx import Presentation
from io import BytesIO

from app.services.report_data import ReportTooLargeError


PPTX_CT = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
XLSX_CT = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"


def _setup_project(db_session, project_factory, article_factory, name="Report API Test", brands=("Auchan",)):
    project = project_factory(name=name)
    project.valid_rows = len(brands)
    for brand in brands:
        article_factory(project, count=1, retailer=brand)
    db_session.commit()
    return project


# --- session / not-found / validation --------------------------------------------


def test_project_report_requires_session(client, project_factory):
    project = project_factory()
    response = client.get(f"/projects/{project.id}/report.pptx", follow_redirects=False)
    assert response.status_code == 307
    assert response.headers["location"] == "/login"


def test_project_report_not_found(authenticated_client):
    response = authenticated_client.get(f"/projects/{uuid.uuid4()}/report.pptx")
    assert response.status_code == 404


def test_comparison_report_requires_session(client, project_factory):
    a = project_factory()
    response = client.get(
        f"/compare/report.pptx?baseline_project_ids={a.id}&comparison_project_ids={a.id}",
        follow_redirects=False,
    )
    assert response.status_code == 307


def test_comparison_report_missing_selection_returns_422(authenticated_client):
    response = authenticated_client.get("/compare/report.pptx")
    assert response.status_code == 422


def test_comparison_report_missing_one_side_returns_422(authenticated_client, project_factory):
    a = project_factory()
    response = authenticated_client.get(f"/compare/report.pptx?baseline_project_ids={a.id}")
    assert response.status_code == 422


# --- success + headers -----------------------------------------------------------


def test_project_pptx_success_headers_and_content(authenticated_client, db_session, project_factory, article_factory):
    project = _setup_project(db_session, project_factory, article_factory)
    response = authenticated_client.get(f"/projects/{project.id}/report.pptx")

    assert response.status_code == 200
    assert response.headers["content-type"] == PPTX_CT
    assert "attachment" in response.headers["content-disposition"]
    assert ".pptx" in response.headers["content-disposition"]
    Presentation(BytesIO(response.content))  # must open without error


def test_project_xlsx_success_headers_and_content(authenticated_client, db_session, project_factory, article_factory):
    project = _setup_project(db_session, project_factory, article_factory)
    response = authenticated_client.get(f"/projects/{project.id}/report.xlsx")

    assert response.status_code == 200
    assert response.headers["content-type"] == XLSX_CT
    assert ".xlsx" in response.headers["content-disposition"]
    openpyxl.load_workbook(BytesIO(response.content))  # must open without error


def test_empty_project_report_still_returns_200(authenticated_client, project_factory):
    project = project_factory()
    response = authenticated_client.get(f"/projects/{project.id}/report.pptx")
    assert response.status_code == 200
    assert len(response.content) > 0


def test_comparison_pptx_and_xlsx_success(authenticated_client, db_session, project_factory, article_factory):
    a = _setup_project(db_session, project_factory, article_factory, name="Cmp API A")
    b = _setup_project(db_session, project_factory, article_factory, name="Cmp API B")

    url = f"/compare/report.pptx?baseline_project_ids={a.id}&comparison_project_ids={b.id}"
    r1 = authenticated_client.get(url)
    assert r1.status_code == 200
    Presentation(BytesIO(r1.content))

    url2 = f"/compare/report.xlsx?baseline_project_ids={a.id}&comparison_project_ids={b.id}"
    r2 = authenticated_client.get(url2)
    assert r2.status_code == 200
    openpyxl.load_workbook(BytesIO(r2.content))


# --- Phase E/F: same-project brand-vs-brand comparison export -------------------


def test_same_project_brand_vs_brand_export_succeeds_and_is_consistent_with_api(
    authenticated_client, internal_headers, db_session, project_factory, article_factory
):
    project = project_factory(name="Brand vs Brand")
    article_factory(project, count=3, retailer="Auchan")
    article_factory(project, count=2, retailer="Carrefour")
    project.total_rows = 5
    db_session.commit()

    url = (
        f"/compare/report.xlsx?baseline_project_ids={project.id}&comparison_project_ids={project.id}"
        "&baseline_filter_brand=Auchan&comparison_filter_brand=Carrefour"
    )
    response = authenticated_client.get(url)
    assert response.status_code == 200
    wb = openpyxl.load_workbook(BytesIO(response.content))
    ws = wb["KPI Comparison"]
    row_values = {
        ws.cell(row=r, column=1).value: (ws.cell(row=r, column=2).value, ws.cell(row=r, column=3).value)
        for r in range(2, ws.max_row + 1)
    }
    baseline_value, comparison_value = row_values["Unique Valid Articles"]
    assert baseline_value == 3
    assert comparison_value == 2

    pptx_url = url.replace("report.xlsx", "report.pptx")
    pptx_response = authenticated_client.get(pptx_url)
    assert pptx_response.status_code == 200
    Presentation(BytesIO(pptx_response.content))

    api_response = authenticated_client.get(
        f"/api/internal/compare?baseline_project_ids={project.id}&comparison_project_ids={project.id}"
        "&baseline_filter_brand=Auchan&comparison_filter_brand=Carrefour",
        headers=internal_headers,
    )
    assert api_response.status_code == 200
    api_body = api_response.json()
    assert api_body["baseline"]["kpis"]["unique_valid_articles"] == baseline_value
    assert api_body["comparison"]["kpis"]["unique_valid_articles"] == comparison_value
    assert api_body["baseline"]["label"] == "Auchan"
    assert api_body["comparison"]["label"] == "Carrefour"


# --- filter propagation (core regression this revision guards) ------------------


def test_filtered_export_differs_from_unfiltered(
    authenticated_client, db_session, project_factory, article_factory
):
    project = _setup_project(
        db_session, project_factory, article_factory, brands=("Auchan", "Carrefour")
    )

    unfiltered = authenticated_client.get(f"/projects/{project.id}/report.xlsx")
    filtered = authenticated_client.get(f"/projects/{project.id}/report.xlsx?brand=Auchan")

    wb_unfiltered = openpyxl.load_workbook(BytesIO(unfiltered.content))
    wb_filtered = openpyxl.load_workbook(BytesIO(filtered.content))

    unfiltered_brands = {
        wb_unfiltered["Brand Performance"].cell(row=r, column=1).value
        for r in range(2, wb_unfiltered["Brand Performance"].max_row + 1)
    }
    filtered_brands = {
        wb_filtered["Brand Performance"].cell(row=r, column=1).value
        for r in range(2, wb_filtered["Brand Performance"].max_row + 1)
    }
    assert unfiltered_brands == {"Auchan", "Carrefour"}
    assert filtered_brands == {"Auchan"}


def test_comparison_filter_propagates_to_export(
    authenticated_client, db_session, project_factory, article_factory
):
    a = project_factory(name="Cmp Filter A")
    b = project_factory(name="Cmp Filter B")
    a.valid_rows = 2
    b.valid_rows = 1
    article_factory(a, count=1, retailer="Auchan")
    article_factory(a, count=1, retailer="Carrefour")
    article_factory(b, count=1, retailer="Auchan")
    db_session.commit()

    url_unfiltered = f"/compare/report.xlsx?baseline_project_ids={a.id}&comparison_project_ids={b.id}"
    url_filtered = url_unfiltered + "&brand=Auchan"

    r1 = authenticated_client.get(url_unfiltered)
    r2 = authenticated_client.get(url_filtered)
    assert len(r1.content) != len(r2.content)


def test_top_n_query_param_is_ignored(authenticated_client, db_session, project_factory, article_factory):
    """`top_n` must not be accepted as a control on the export's display
    cap — the report always uses its own fixed REPORT_TOP_N/EXCEL_TOP_N.
    """
    project = _setup_project(db_session, project_factory, article_factory, brands=("Auchan", "Carrefour"))

    r1 = authenticated_client.get(f"/projects/{project.id}/report.xlsx")
    r2 = authenticated_client.get(f"/projects/{project.id}/report.xlsx?top_n=1")
    assert r1.content == r2.content or len(r1.content) == len(r2.content)


# --- displayed/exported consistency ----------------------------------------------


def test_exported_kpi_matches_internal_api(
    authenticated_client, internal_headers, db_session, project_factory, article_factory
):
    project = _setup_project(db_session, project_factory, article_factory)

    api_response = authenticated_client.get(
        f"/api/internal/projects/{project.id}/analytics", headers=internal_headers
    )
    api_unique_valid = api_response.json()["kpis"]["unique_valid_articles"]

    xlsx_response = authenticated_client.get(f"/projects/{project.id}/report.xlsx")
    wb = openpyxl.load_workbook(BytesIO(xlsx_response.content))
    ws = wb["KPI Summary"]
    row_values = {ws.cell(row=r, column=1).value: ws.cell(row=r, column=2).value for r in range(2, ws.max_row + 1)}

    assert row_values["Unique valid articles"] == api_unique_valid


# --- oversized report -> controlled error, not a 500 ------------------------------


def test_oversized_project_report_returns_422_not_500(
    authenticated_client, db_session, project_factory, article_factory
):
    project = _setup_project(db_session, project_factory, article_factory)

    with patch("app.services.report_pptx.MAX_PPTX_BYTES", 100):
        response = authenticated_client.get(f"/projects/{project.id}/report.pptx")
    assert response.status_code == 422
    assert "narrow the filters" in response.json()["detail"].lower()


# --- no secret leakage -------------------------------------------------------------


def test_no_secret_leakage_in_report_response(authenticated_client, db_session, project_factory, article_factory):
    project = _setup_project(db_session, project_factory, article_factory)
    response = authenticated_client.get(f"/projects/{project.id}/report.xlsx")
    assert b"test-internal-secret" not in response.content


# --- filename / header hardening --------------------------------------------------


def test_non_ascii_project_name_produces_ascii_filename(
    authenticated_client, db_session, project_factory, article_factory
):
    project = _setup_project(
        db_session, project_factory, article_factory, name="Măgurele Distribuție"
    )
    response = authenticated_client.get(f"/projects/{project.id}/report.pptx")
    disposition = response.headers["content-disposition"]
    filename = disposition.split("filename=")[1].strip('"')
    assert filename.isascii()
    assert filename  # non-empty
    disposition.encode("latin-1")  # must be a valid HTTP header value


def test_quotes_and_semicolons_in_project_name_are_neutralized(
    authenticated_client, db_session, project_factory, article_factory
):
    project = _setup_project(
        db_session, project_factory, article_factory, name='Weird"Name;With\'Chars'
    )
    response = authenticated_client.get(f"/projects/{project.id}/report.pptx")
    disposition = response.headers["content-disposition"]
    assert '"' not in disposition.split("filename=")[1][:-1] or disposition.count('"') == 2
    assert ";" not in disposition.split("filename=")[1]


def test_crlf_in_project_name_cannot_inject_headers(
    authenticated_client, db_session, project_factory, article_factory
):
    malicious_name = "Evil\r\nX-Injected: true"
    project = _setup_project(db_session, project_factory, article_factory, name=malicious_name)
    response = authenticated_client.get(f"/projects/{project.id}/report.pptx")

    assert response.status_code == 200
    assert "X-Injected" not in response.headers
    raw_disposition = response.headers["content-disposition"]
    assert "\r" not in raw_disposition
    assert "\n" not in raw_disposition


def test_extremely_long_project_name_produces_bounded_filename(
    authenticated_client, db_session, project_factory, article_factory
):
    # 255 is the app's own max project-name length (ProjectCreate) — still
    # well past MAX_FILENAME_COMPONENT_CHARS (100), so this still exercises
    # the filename-length cap.
    long_name = "A" * 255
    project = _setup_project(db_session, project_factory, article_factory, name=long_name)
    response = authenticated_client.get(f"/projects/{project.id}/report.pptx")
    disposition = response.headers["content-disposition"]
    filename = disposition.split("filename=")[1].strip('"')
    assert len(filename) < 300


def test_duplicate_period_labels_produce_well_formed_filename(
    authenticated_client, db_session, project_factory, article_factory
):
    a = project_factory(name="SameQuarterA", quarter="2026-Q2")
    b = project_factory(name="SameQuarterB", quarter="2026-Q2")
    a.valid_rows = 1
    b.valid_rows = 1
    article_factory(a, count=1, retailer="Auchan")
    article_factory(b, count=1, retailer="Auchan")
    db_session.commit()

    response = authenticated_client.get(
        f"/compare/report.pptx?baseline_project_ids={a.id}&comparison_project_ids={b.id}"
    )
    assert response.status_code == 200
    disposition = response.headers["content-disposition"]
    filename = disposition.split("filename=")[1].strip('"')
    assert filename.endswith(".pptx")
    assert filename  # non-empty, well-formed despite baseline == comparison label

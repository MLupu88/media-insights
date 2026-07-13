from io import BytesIO
from unittest.mock import patch

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor

from app.services import report_pptx
from app.services.analytics import AnalyticsFilters
from app.services.narrative_service import create_project_generation
from app.services.report_contract import REPORT_TOP_N
from app.services.report_data import (
    ReportTooLargeError,
    build_comparison_report_data,
    build_project_report_data,
)


def _slide_texts(prs) -> list[list[str]]:
    result = []
    for slide in prs.slides:
        texts = [
            shape.text_frame.text for shape in slide.shapes if shape.has_text_frame and shape.text_frame.text
        ]
        result.append(texts)
    return result


def _open(data: bytes) -> Presentation:
    return Presentation(BytesIO(data))


def test_project_pptx_slide_structure(report_db, project_factory, article_factory):
    project = project_factory()
    for i in range(15):
        article_factory(project, count=1, retailer=f"Brand{i}")

    data = build_project_report_data(report_db, project.id, AnalyticsFilters())
    pptx_bytes = report_pptx.build_project_pptx(data)
    prs = _open(pptx_bytes)

    texts = _slide_texts(prs)
    headings = [t[0].split("\n")[0] for t in texts if t]
    assert headings == [
        project.name,
        "Executive Summary",
        "Brand & Competitor Performance",
        "Topic & Category Mix",
        "Sentiment & Brand Role",
        "Publications",
        "Story Clusters",
        "Methodology & Filters",
    ]


def test_project_pptx_empty_variant_is_three_slides(report_db, project_factory):
    project = project_factory()
    data = build_project_report_data(report_db, project.id, AnalyticsFilters())
    prs = _open(report_pptx.build_project_pptx(data))

    assert len(prs.slides) == 3
    texts = _slide_texts(prs)
    assert texts[1][0] == "No Data Available"
    assert texts[2][0].startswith("Methodology")


def test_comparison_pptx_slide_structure(report_db, project_factory, article_factory):
    a = project_factory(name="A", quarter="2026-Q1")
    b = project_factory(name="B", quarter="2026-Q2")
    article_factory(a, count=1, retailer="Auchan")
    article_factory(b, count=1, retailer="Auchan")

    data = build_comparison_report_data(report_db, [a.id], [b.id], AnalyticsFilters())
    prs = _open(report_pptx.build_comparison_pptx(data))

    texts = _slide_texts(prs)
    headings = [t[0].split("\n")[0] for t in texts if t]
    assert headings == [
        "Comparison Report",
        "Executive Summary",
        "Brand & SOV Movement",
        "Topic & Category Shifts",
        "Sentiment & Brand-Role Change",
        "Publication Movement",
        "Story Movement",
        "Volatility & Concentration",
        "Methodology & Filters",
    ]


def test_comparison_pptx_empty_variant_is_three_slides(report_db, project_factory):
    a = project_factory(name="EmptyA", quarter="2026-Q1")
    b = project_factory(name="EmptyB", quarter="2026-Q2")
    data = build_comparison_report_data(report_db, [a.id], [b.id], AnalyticsFilters())
    prs = _open(report_pptx.build_comparison_pptx(data))

    assert len(prs.slides) == 3
    texts = _slide_texts(prs)
    assert texts[1][0] == "No Data Available"


# --- legibility: capping, captions, label truncation -----------------------------


def test_brand_slide_capped_at_report_top_n_with_caption(report_db, project_factory, article_factory):
    project = project_factory()
    for i in range(15):
        article_factory(project, count=1, retailer=f"Brand{i}")

    data = build_project_report_data(report_db, project.id, AnalyticsFilters())
    prs = _open(report_pptx.build_project_pptx(data))

    brand_slide = prs.slides[2]
    table_shapes = [s for s in brand_slide.shapes if s.has_table]
    assert len(table_shapes) == 1
    table = table_shapes[0].table
    assert len(table.rows) - 1 == REPORT_TOP_N  # minus header row

    texts = [s.text_frame.text for s in brand_slide.shapes if s.has_text_frame]
    assert any(f"Top {REPORT_TOP_N} of 15" in t for t in texts)


def test_no_caption_when_under_cap(report_db, project_factory, article_factory):
    project = project_factory()
    article_factory(project, count=1, retailer="OnlyBrand")

    data = build_project_report_data(report_db, project.id, AnalyticsFilters())
    prs = _open(report_pptx.build_project_pptx(data))

    brand_slide = prs.slides[2]
    texts = [s.text_frame.text for s in brand_slide.shapes if s.has_text_frame]
    assert not any("shown" in t for t in texts)


def test_long_label_truncated_on_slide(report_db, project_factory, article_factory):
    long_name = "A" * 60
    project = project_factory()
    article_factory(project, count=1, retailer=long_name)

    data = build_project_report_data(report_db, project.id, AnalyticsFilters())
    prs = _open(report_pptx.build_project_pptx(data))

    brand_slide = prs.slides[2]
    table = next(s for s in brand_slide.shapes if s.has_table).table
    cell_text = table.cell(1, 0).text
    assert cell_text != long_name
    assert len(cell_text) <= 28
    assert cell_text.endswith("…")


# --- brand/accent styling ----------------------------------------------------------


def test_accent_color_applied_on_section_slide(report_db, project_factory, article_factory):
    from pptx.enum.dml import MSO_FILL_TYPE

    project = project_factory()
    article_factory(project, count=1, retailer="Auchan")

    data = build_project_report_data(report_db, project.id, AnalyticsFilters())
    prs = _open(report_pptx.build_project_pptx(data))

    brand_slide = prs.slides[2]
    accent_shapes = []
    for shape in brand_slide.shapes:
        if not hasattr(shape, "fill"):
            continue
        if shape.fill.type != MSO_FILL_TYPE.SOLID:
            continue
        if shape.fill.fore_color.rgb == RGBColor(0xB3, 0x49, 0x2B):
            accent_shapes.append(shape)
    assert len(accent_shapes) >= 1


# --- insight rendering -----------------------------------------------------------


def test_valid_insight_rendered_with_label(report_db, db_session, project_factory, article_factory):
    project = project_factory()
    article_factory(project, count=1, retailer="Auchan")
    generation, _ = create_project_generation(db_session, project, AnalyticsFilters())
    generation.status = "complete"
    db_session.commit()

    from app.models.narrative import NarrativeInsight, NarrativeValidationStatus

    insight = NarrativeInsight(
        generation_id=generation.id, project_id=generation.project_id, narrative_type="recommendations",
        key="r1", title="Increase digital coverage", narrative="Consider expanding digital outreach.",
        evidence_type="kpi_delta", evidence=[], raw_candidate={},
        validation_status=NarrativeValidationStatus.VALID,
    )
    db_session.add(insight)
    db_session.commit()

    data = build_project_report_data(report_db, project.id, AnalyticsFilters())
    prs = _open(report_pptx.build_project_pptx(data))

    all_text = "\n".join(t for slide_texts in _slide_texts(prs) for t in slide_texts)
    assert "Increase digital coverage" in all_text
    assert "RECOMMENDATION" in all_text


def test_rejected_insight_never_rendered(report_db, db_session, project_factory, article_factory):
    project = project_factory()
    article_factory(project, count=1, retailer="Auchan")
    generation, _ = create_project_generation(db_session, project, AnalyticsFilters())
    generation.status = "complete"
    db_session.commit()

    from app.models.narrative import NarrativeInsight, NarrativeValidationStatus

    insight = NarrativeInsight(
        generation_id=generation.id, project_id=generation.project_id, narrative_type="key_findings",
        key="k1", title="Should Never Appear In Deck", narrative="Rejected content.",
        evidence_type="kpi_delta", evidence=[], raw_candidate={},
        validation_status=NarrativeValidationStatus.REJECTED, rejection_reason="bad",
    )
    db_session.add(insight)
    db_session.commit()

    data = build_project_report_data(report_db, project.id, AnalyticsFilters())
    prs = _open(report_pptx.build_project_pptx(data))

    all_text = "\n".join(t for slide_texts in _slide_texts(prs) for t in slide_texts)
    assert "Should Never Appear In Deck" not in all_text


# --- methodology slide -------------------------------------------------------------


def test_methodology_slide_contains_required_fields(report_db, project_factory, article_factory):
    project = project_factory()
    article_factory(project, count=1, retailer="Auchan")

    data = build_project_report_data(report_db, project.id, AnalyticsFilters())
    prs = _open(report_pptx.build_project_pptx(data))

    methodology_text = _slide_texts(prs)[-1]
    combined = "\n".join(methodology_text)
    assert "Scope:" in combined
    assert "Filters:" in combined
    assert "Generated:" in combined
    assert "UTC" in combined
    assert "unique valid articles" in combined
    assert "Article detail" in combined
    assert "Insights:" in combined


# --- table styling: no black-rectangle regression --------------------------------


def test_table_body_cells_do_not_use_header_fill(report_db, project_factory, article_factory):
    """The bug: a sparse table's body rows inherited the header's solid
    INK fill because PowerPoint stretched a large caller-supplied height
    evenly across too few rows. Header must stay INK/black; body rows must
    be WHITE/zebra, never INK.
    """
    from pptx.enum.dml import MSO_FILL_TYPE

    project = project_factory()
    article_factory(project, count=1, retailer="Auchan")

    data = build_project_report_data(report_db, project.id, AnalyticsFilters())
    prs = _open(report_pptx.build_project_pptx(data))

    brand_slide = prs.slides[2]
    table = next(s for s in brand_slide.shapes if s.has_table).table

    header_cell = table.cell(0, 0)
    assert header_cell.fill.type == MSO_FILL_TYPE.SOLID
    assert header_cell.fill.fore_color.rgb == RGBColor(0x11, 0x11, 0x11)

    body_cell = table.cell(1, 0)
    assert body_cell.fill.type == MSO_FILL_TYPE.SOLID
    assert body_cell.fill.fore_color.rgb != RGBColor(0x11, 0x11, 0x11)
    assert body_cell.fill.fore_color.rgb == RGBColor(0xFF, 0xFF, 0xFF)


def test_table_row_height_is_not_stretched_for_sparse_tables(report_db, project_factory, article_factory):
    """A 2-row table (header + 1 data row) must not be evenly stretched
    across a large fixed height — each row stays a compact, fixed height.
    """
    from pptx.util import Emu

    project = project_factory()
    article_factory(project, count=1, retailer="Auchan")

    data = build_project_report_data(report_db, project.id, AnalyticsFilters())
    prs = _open(report_pptx.build_project_pptx(data))

    brand_slide = prs.slides[2]
    table = next(s for s in brand_slide.shapes if s.has_table).table

    for row in table.rows:
        assert row.height < Emu(500000)  # well under half an inch; never a stretched black block


def test_empty_section_renders_no_data_message_not_black_table(report_db, project_factory, article_factory):
    """A project with articles but no classification leaves the Topic
    slide's ranked section empty — it must render a compact "No data
    available" message, never an empty/oversized table shape.
    """
    project = project_factory()
    article_factory(project, count=1, retailer="Auchan")  # unclassified: topic distribution is empty

    data = build_project_report_data(report_db, project.id, AnalyticsFilters())
    prs = _open(report_pptx.build_project_pptx(data))

    topic_slide = prs.slides[3]
    assert not any(s.has_table for s in topic_slide.shapes)
    texts = [s.text_frame.text for s in topic_slide.shapes if s.has_text_frame]
    assert any("No data available" in t for t in texts)


# --- size bound ----------------------------------------------------------------


def test_oversized_pptx_raises_controlled_error(report_db, project_factory, article_factory):
    project = project_factory()
    article_factory(project, count=1, retailer="Auchan")
    data = build_project_report_data(report_db, project.id, AnalyticsFilters())

    with patch("app.services.report_pptx.MAX_PPTX_BYTES", 100):
        with pytest.raises(ReportTooLargeError):
            report_pptx.build_project_pptx(data)
